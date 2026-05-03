from __future__ import annotations

import math
from copy import deepcopy
from typing import Sequence
from pptx.util import Pt

from PIL import Image

from .dependencies import require_dependency
from .models import FeatureRow, ListingExtraction, RoomSection, SelectedImage
from .utils import format_beds_and_baths, format_total, normalize_sqft_text


def _pt(value: float):
    pptx_util = require_dependency("pptx.util", "python-pptx")
    return pptx_util.Pt(value)


def build_text_mapping(listing: ListingExtraction, brochure_copy, gallery_images: Sequence[SelectedImage]) -> dict[str, str]:
    caption_map = {item.file_name: item.caption for item in brochure_copy.captions}
    mapping = {
        "{{ADDRESS}}": listing.short_address,
        "{{PRICE}}": listing.price,
        "{{BEDS_AND_BATHS}}": format_beds_and_baths(listing.bedrooms, listing.bathrooms),
        "{{TITLE}}": brochure_copy.title,
        "{{TOTAL}}": format_total(listing.total_sqft),
    }
    for index in range(1, 11):
        caption = ""
        if index <= len(gallery_images):
            image = gallery_images[index - 1]
            caption = caption_map.get(image.path.name, image.caption)
        mapping[f"{{{{caption_{index}}}}}"] = caption
    return mapping


def clone_run_style(source_run, dest_run) -> None:
    if source_run is None:
        return
    font = source_run.font
    dest_font = dest_run.font
    dest_font.bold = font.bold
    dest_font.italic = font.italic
    dest_font.name = font.name
    dest_font.size = font.size
    if font.color is not None and getattr(font.color, "rgb", None) is not None:
        dest_font.color.rgb = font.color.rgb


def apply_paragraph_style(source, dest) -> None:
    dest.alignment = source.alignment
    dest.level = source.level
    dest.line_spacing = source.line_spacing
    dest.space_before = source.space_before
    dest.space_after = source.space_after
    source_ppr = getattr(source._p, "pPr", None)
    if source_ppr is not None:
        dest_ppr = getattr(dest._p, "pPr", None)
        if dest_ppr is not None:
            dest._p.remove(dest_ppr)
        dest._p.insert(0, deepcopy(source_ppr))


def capture_first_run_style(paragraph) -> dict:
    run = paragraph.runs[0] if paragraph.runs else None
    if run is None:
        return {}
    color = None
    if run.font.color is not None and getattr(run.font.color, "rgb", None) is not None:
        color = run.font.color.rgb
    return {
        "bold": run.font.bold,
        "italic": run.font.italic,
        "name": run.font.name,
        "size": run.font.size,
        "color": color,
    }


def apply_run_style_from_snapshot(run, snapshot: dict) -> None:
    font = run.font
    font.bold = snapshot.get("bold")
    font.italic = snapshot.get("italic")
    font.name = snapshot.get("name")
    font.size = snapshot.get("size")
    color = snapshot.get("color")
    if color is not None:
        font.color.rgb = color


def replace_text_in_shape(shape, replacements: dict[str, str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        original = "".join(run.text for run in paragraph.runs)
        if not original:
            continue
        updated = original
        for placeholder, replacement in replacements.items():
            updated = updated.replace(placeholder, replacement)
        if updated == original:
            continue
        source_run = paragraph.runs[0] if paragraph.runs else None
        alignment = paragraph.alignment
        level = paragraph.level
        paragraph.text = updated
        paragraph.alignment = alignment
        paragraph.level = level
        if paragraph.runs:
            clone_run_style(source_run, paragraph.runs[0])


def populate_summary_block(shape, bullets: Sequence[str]) -> None:
    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()
    values = list(bullets)[:6] or [""]
    for index, bullet in enumerate(values):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        apply_paragraph_style(template_paragraph, paragraph)
        paragraph.text = bullet
        if paragraph.runs:
            apply_run_style_from_snapshot(paragraph.runs[0], run_style)
    distribute_summary_spacing(shape, text_frame, run_style)


def distribute_summary_spacing(shape, text_frame, run_style: dict) -> None:
    paragraphs = [p for p in text_frame.paragraphs if "".join(run.text for run in p.runs).strip()]
    if len(paragraphs) < 2:
        return

    font_size = run_style.get("size")
    font_points = font_size.pt if font_size is not None else 14
    chars_per_line = max(24, int(shape.width.pt / (font_points * 0.48)))
    visual_line_count = 0

    for paragraph in paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        visual_line_count += max(1, math.ceil(len(text) / chars_per_line))

    estimated_text_height = visual_line_count * font_points * 1.15
    remaining = max(0.0, shape.height.pt - estimated_text_height)
    per_gap = remaining / max(1, len(paragraphs) - 1)

    # Leave a little safety because PowerPoint wrapping is not exposed by python-pptx.
    per_gap = max(2.0, min(18.0, per_gap * 0.9))
    for idx, paragraph in enumerate(paragraphs):
        paragraph.space_before = _pt(0 if idx == 0 else per_gap)
        paragraph.space_after = _pt(0)
        paragraph.line_spacing = 1.0


def format_feature_parent_line(label: str, value: str) -> str:
    tabs = "\t" if len(label) <= 5 else "\t"
    return f"{label}:{tabs}{value}"


def format_feature_child_line(value: str) -> str:
    return f"\t\t{value}"


def populate_feature_block(shape, feature_rows: Sequence[FeatureRow]) -> None:
    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()

    line_types: list[str] = []
    first = True
    for row in feature_rows:
        values = [value.strip() for value in row.values if value.strip()]
        if not values:
            continue
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        apply_paragraph_style(template_paragraph, paragraph)
        paragraph.text = format_feature_parent_line(row.label, values[0])
        line_types.append("parent")
        for run in paragraph.runs:
            apply_run_style_from_snapshot(run, run_style)

        for extra_value in values[1:]:
            extra = text_frame.add_paragraph()
            apply_paragraph_style(template_paragraph, extra)
            extra.text = format_feature_child_line(extra_value)
            line_types.append("child")
            for run in extra.runs:
                apply_run_style_from_snapshot(run, run_style)

    distribute_feature_spacing(shape, text_frame, line_types, run_style)


def distribute_feature_spacing(shape, text_frame, line_types: Sequence[str], run_style: dict) -> None:
    paragraphs = [p for p in text_frame.paragraphs if "".join(run.text for run in p.runs).strip()]
    if len(paragraphs) <= 1:
        return

    font_size = run_style.get("size")
    font_pt = font_size.pt if font_size else 11
    
    # 1. Calculate the 'Floor' and 'Ceiling' of your spacing needs in points
    # We are calculating the extra gap needed ABOVE each line (space_before)
    # Child range: 1.0 to 1.5 total (so 0.0*font to 0.5*font extra)
    # Parent range: 2.0 to 2.5 total (so 1.0*font to 1.5*font extra)
    
    min_needed_gap = 0
    max_needed_gap = 0
    
    for i in range(1, len(paragraphs)):
        l_type = line_types[i] if i < len(line_types) else "child"
        if l_type == "parent":
            min_needed_gap += 1.0 * font_pt
            max_needed_gap += 1.5 * font_pt
        else:
            min_needed_gap += 0.0 * font_pt
            max_needed_gap += 0.5 * font_pt
            
    # Calculate physical space available for GAPS 
    # (Box Height - space taken by the actual text lines)
    # Using 1.1 as a tight base multiplier for the text itself
    available_gap_space = shape.height.pt - (len(paragraphs) * font_pt * 1.1)
    
    # 2. Determine the 'Fill Factor' (k)
    # This finds where we land between our Min and Max targets
    if max_needed_gap > min_needed_gap:
        k = (available_gap_space - min_needed_gap) / (max_needed_gap - min_needed_gap)
    else:
        k = 0
    
    # Clamp k between 0 and 1 so we never exceed your 1.5/2.5 limits
    k = max(0.0, min(1.0, k))

    # 3. Apply the Interpolated Spacing
    for i, paragraph in enumerate(paragraphs):
        if i == 0:
            paragraph.space_before = Pt(0)
            continue
            
        l_type = line_types[i] if i < len(line_types) else "child"
        
        if l_type == "parent":
            # Map k to the 1.0 -> 1.5 range
            extra_space = font_pt * (1.0 + (k * 0.5))
        else:
            # Map k to the 0.0 -> 0.5 range
            extra_space = font_pt * (0.0 + (k * 0.5))
            
        paragraph.space_before = Pt(extra_space)
        paragraph.space_after = Pt(0)
        # Force single line spacing so our space_before math is the only driver
        paragraph.line_spacing = 1.0

def populate_room_block(shape, room_sections: Sequence[RoomSection]) -> None:
    from pptx.enum.text import PP_ALIGN

    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()

    first = True
    for section in room_sections:
        heading = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        apply_paragraph_style(template_paragraph, heading)
        heading.alignment = PP_ALIGN.CENTER
        heading.text = f"{section.title} ({normalize_sqft_text(section.area_text)})"
        for run in heading.runs:
            apply_run_style_from_snapshot(run, run_style)
            run.font.italic = True

        spacer = text_frame.add_paragraph()
        spacer.text = ""

        for room in section.rooms:
            row = text_frame.add_paragraph()
            apply_paragraph_style(template_paragraph, row)
            row.alignment = None
            row.text = f"{room.name}\t{room.size}"
            for run in row.runs:
                apply_run_style_from_snapshot(run, run_style)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def fill_picture_placeholder(slide, placeholder_shape, image_path):
    left, top, width, height = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    with Image.open(image_path) as img:
        image_width, image_height = img.size

    picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio > box_ratio:
        scaled_width = height * image_ratio
        crop = (scaled_width - width) / scaled_width / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        scaled_height = width / image_ratio
        crop = (scaled_height - height) / scaled_height / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    remove_shape(placeholder_shape)


def render_presentation(template_path, output_path, text_mapping, feature_rows, room_sections, main_image, gallery_images) -> None:
    pptx = require_dependency("pptx", "python-pptx")
    presentation = pptx.Presentation(str(template_path))
    image_map = {"{{IMG_MAIN}}": main_image.path}
    for index, image in enumerate(gallery_images, start=1):
        image_map[f"{{{{IMG_{index}}}}}"] = image.path

    for slide in presentation.slides:
        for shape in list(slide.shapes):
            shape_name = getattr(shape, "name", "").strip()
            if shape_name in image_map:
                fill_picture_placeholder(slide, shape, image_map[shape_name])
                continue
            if getattr(shape, "has_text_frame", False):
                combined_text = "\n".join("".join(run.text for run in paragraph.runs) for paragraph in shape.text_frame.paragraphs)
                if "{{SUMMARY}}" in combined_text:
                    populate_summary_block(shape, text_mapping.get("{{SUMMARY}}", "").split("\n"))
                elif "{{FEATURE_BLOCK}}" in combined_text:
                    populate_feature_block(shape, feature_rows)
                elif "{{ROOM_BLOCK}}" in combined_text:
                    populate_room_block(shape, room_sections)
                else:
                    replace_text_in_shape(shape, text_mapping)
    presentation.save(str(output_path))
