from __future__ import annotations

import argparse
import base64
import io
import os
import re
import sys
import threading
import traceback
from collections import Counter
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Sequence

from openai import OpenAI
from pydantic import BaseModel, Field
from PIL import Image
from pptx.util import Pt


SUPPORTED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
DEFAULT_MODEL = "gpt-5.4-nano"


class FeatureRow(BaseModel):
    label: str
    values: list[str] = Field(default_factory=list)


class RoomEntry(BaseModel):
    name: str
    size: str


class RoomSection(BaseModel):
    title: str
    area_text: str
    rooms: list[RoomEntry] = Field(default_factory=list)


class ListingExtraction(BaseModel):
    address: str
    short_address: str
    price: str
    bedrooms: int | None = None
    bathrooms: int | None = None
    total_sqft: str | None = None
    summary_bullets: list[str] = Field(default_factory=list)
    feature_rows: list[FeatureRow] = Field(default_factory=list)
    room_sections: list[RoomSection] = Field(default_factory=list)


class ImageInsight(BaseModel):
    file_name: str
    scene_type: str
    room_type: str
    caption: str
    showcase_score: int = Field(ge=1, le=10)
    front_exterior: bool = False
    exterior: bool = False
    avoid: bool = False
    avoid_reason: str = ""


class ImageBatchAnalysis(BaseModel):
    images: list[ImageInsight]


class TitleResponse(BaseModel):
    title: str


class BrochureCopyResponse(BaseModel):
    title: str
    summary_bullets: list[str] = Field(default_factory=list)
    captions: list["CaptionItem"] = Field(default_factory=list)


class CaptionItem(BaseModel):
    file_name: str
    caption: str


@dataclass
class SelectedImage:
    path: Path
    caption: str
    showcase_score: int
    scene_type: str
    room_type: str
    front_exterior: bool
    exterior: bool


@dataclass
class AppConfig:
    pdf_path: Path
    image_dir: Path
    template_path: Path
    output_path: Path
    model: str = DEFAULT_MODEL


def log_print(message: str) -> None:
    print(message)


def require_dependency(module_name: str, package_name: str):
    try:
        return __import__(module_name, fromlist=["*"])
    except ImportError as exc:
        raise RuntimeError(
            f"Missing dependency '{package_name}'. Install the packages from requirements.txt first."
        ) from exc


def build_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is not set.")

    base_url = os.getenv("OPENAI_BASE_URL")
    if base_url:
        return OpenAI(api_key=api_key, base_url=base_url)
    return OpenAI(api_key=api_key)


def extract_pdf_text(pdf_path: Path) -> str:
    pypdf = require_dependency("pypdf", "pypdf")
    reader = pypdf.PdfReader(str(pdf_path))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    text = "\n".join(pages).strip()
    if not text:
        raise RuntimeError(f"No text could be extracted from {pdf_path}.")
    return text


def analyze_listing_text(client: OpenAI, model: str, raw_text: str) -> ListingExtraction:
    instructions = (
        "You extract real-estate brochure data from messy listing text. "
        "Return only brochure-safe factual data. "
        "Do not invent facts. "
        "Create a short_address that is just the street address unit and street, like '4-1616 Happyvale Ave'. "
        "Do not include city, province, or postal code in short_address. "
        "Create 6-7 summary_bullets that each read like a strong brochure bullet point, not a sentence fragment. "
        "Each bullet should be roughly 18-32 words. "
        "Normalize labels and keep room order as presented. "
        "For feature_rows, output rows in this preferred order when available: "
        "Location, Style, Heating, Fireplace, Parking, Includes, Age, Taxes, Lot Size. "
        "Each row should contain one label and concise brochure-friendly values. "
        "Examples: Location='Brocklehurst'; Style='Cathedral Entry'; Heating='Forced Air' and 'Central Air'; "
        "Parking='2 Car Garage' and 'Additional Parking'; Includes='Dishwasher', 'Electric Range', 'Refrigerator', 'Washer/Dryer'; "
        "Age='42 Years'; Taxes='$3,145 (2025)'; Lot Size='.23 Acres'. "
        "Do not include MLS numbers, strata fee details, major area labels, or long administrative text in feature_rows. "
        "For room_sections, preserve meaningful sections like Main/Basement/Upstairs/Lower where possible. "
        "Each section title should be short and brochure-ready, and area_text should be a concise value like '1,081 sq. ft.' "
        "without repeating the word Total. "
        "Use area hints from nearby values such as Above Grade, Below Grade, 1st, Lower, or Bsmt when present. "
        "Bathrooms may use values like '4 piece' or '2 piece'. "
        "If an item is unavailable, omit it or leave it null."
    )

    response = client.responses.parse(
        model=model,
        instructions=instructions,
        input=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_text",
                        "text": (
                            "Extract brochure data from this listing text into JSON.\n\n"
                            f"{raw_text}"
                        ),
                    }
                ],
            }
        ],
        text_format=ListingExtraction,
        text={"verbosity": "low"},
        reasoning={"effort": "low"},
        max_output_tokens=4000,
    )
    parsed = response.output_parsed
    if not parsed:
        raise RuntimeError("The model returned no parsed listing data.")
    return parsed


def iter_image_paths(image_dir: Path) -> list[Path]:
    images = [p for p in sorted(image_dir.iterdir()) if p.suffix.lower() in SUPPORTED_IMAGE_EXTENSIONS]
    if not images:
        raise RuntimeError(f"No supported images were found in {image_dir}.")
    return images


def encode_image_for_model(image_path: Path, max_size: int = 1024) -> str:
    with Image.open(image_path) as img:
        image = img.convert("RGB")
        image.thumbnail((max_size, max_size))
        buffer = io.BytesIO()
        image.save(buffer, format="JPEG", quality=72, optimize=True)
        encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/jpeg;base64,{encoded}"


def chunked(items: Sequence[Path], size: int) -> Iterable[Sequence[Path]]:
    for index in range(0, len(items), size):
        yield items[index : index + size]


def analyze_image_batch(client: OpenAI, model: str, image_paths: Sequence[Path]) -> list[ImageInsight]:
    content: list[dict] = [
        {
            "type": "input_text",
            "text": (
                "Review these real-estate listing photos for brochure selection. "
                "For each image, determine if it should be avoided, whether it is the front exterior, "
                "whether it is another exterior, assign a showcase_score from 1-10, classify the room_type, "
                "and write a short brochure caption that highlights the best feature of the photo. "
                "Prefer bright, wide, well-composed photos that help sell the house. "
                "Avoid blurry images, duplicate angles, closeups with little context, floorplans, logos, "
                "screenshots, and utility-only shots unless there are very few choices. "
                "Set scene_type to one of: front_exterior, backyard, deck_patio, living_room, dining_room, kitchen, "
                "primary_bedroom, bedroom, bathroom, rec_room, laundry, garage, view, other_exterior, other."
            ),
        }
    ]

    for image_path in image_paths:
        content.append({"type": "input_text", "text": f"FILE: {image_path.name}"})
        content.append(
            {
                "type": "input_image",
                "image_url": encode_image_for_model(image_path),
                "detail": "low",
            }
        )

    response = client.responses.parse(
        model=model,
        instructions="Return structured JSON only.",
        input=[{"role": "user", "content": content}],
        text_format=ImageBatchAnalysis,
        text={"verbosity": "low"},
        reasoning={"effort": "low"},
        max_output_tokens=2500,
    )
    parsed = response.output_parsed
    if not parsed:
        raise RuntimeError("The model returned no parsed image data.")
    return parsed.images


def analyze_images(
    client: OpenAI, model: str, image_paths: Sequence[Path], logger: Callable[[str], None]
) -> list[SelectedImage]:
    insights: list[ImageInsight] = []
    batches = list(chunked(list(image_paths), 8))
    for idx, batch in enumerate(batches, start=1):
        logger(f"Analyzing image batch {idx}/{len(batches)}...")
        insights.extend(analyze_image_batch(client, model, batch))

    path_by_name = {path.name: path for path in image_paths}
    cleaned: list[SelectedImage] = []
    for insight in insights:
        image_path = path_by_name.get(insight.file_name)
        if not image_path or insight.avoid:
            continue
        cleaned.append(
            SelectedImage(
                path=image_path,
                caption=clean_caption(insight.caption),
                showcase_score=insight.showcase_score,
                scene_type=normalize_scene_type(insight.scene_type),
                room_type=insight.room_type.strip() or "Photo",
                front_exterior=insight.front_exterior,
                exterior=insight.exterior,
            )
        )

    if not cleaned:
        raise RuntimeError("All images were filtered out. Try using a different folder or model.")

    # Keep the highest-scoring image per filename and sort strongest first.
    unique: dict[Path, SelectedImage] = {}
    for image in cleaned:
        existing = unique.get(image.path)
        if existing is None or image.showcase_score > existing.showcase_score:
            unique[image.path] = image
    ranked = sorted(unique.values(), key=lambda item: item.showcase_score, reverse=True)
    return ranked


def clean_caption(text: str) -> str:
    value = " ".join(text.split())
    if not value:
        return "Beautiful space with strong natural light and inviting everyday comfort."
    if len(value) > 92:
        value = value[:92].rsplit(" ", 1)[0]
    value = value.rstrip(". ")
    return f"{value}."


def normalize_scene_type(scene_type: str) -> str:
    value = scene_type.strip().lower().replace(" ", "_")
    allowed = {
        "front_exterior",
        "backyard",
        "deck_patio",
        "living_room",
        "dining_room",
        "kitchen",
        "primary_bedroom",
        "bedroom",
        "bathroom",
        "rec_room",
        "laundry",
        "garage",
        "view",
        "other_exterior",
        "other",
    }
    return value if value in allowed else "other"


def choose_images(ranked_images: Sequence[SelectedImage]) -> tuple[SelectedImage, list[SelectedImage]]:
    main_image = next((img for img in ranked_images if img.front_exterior), None)
    if main_image is None:
        main_image = next((img for img in ranked_images if img.exterior), None)
    if main_image is None:
        main_image = ranked_images[0]

    candidates = [img for img in ranked_images if img.path != main_image.path]
    selected: list[SelectedImage] = []
    scene_counts: Counter[str] = Counter()
    exterior_count = 0

    def exterior_like(image: SelectedImage) -> bool:
        return image.scene_type in {"front_exterior", "backyard", "deck_patio", "view", "other_exterior"}

    for image in candidates:
        if len(selected) >= 10:
            break
        if scene_counts[image.scene_type] >= 1 and len(selected) < 8:
            continue
        if image.scene_type in {"bedroom", "primary_bedroom"} and scene_counts["bedroom"] + scene_counts["primary_bedroom"] >= 2:
            continue
        if exterior_like(image) and exterior_count >= 3 and len(selected) < 9:
            continue
        selected.append(image)
        scene_counts[image.scene_type] += 1
        if exterior_like(image):
            exterior_count += 1

    for image in candidates:
        if len(selected) >= 10:
            break
        if image.path in {item.path for item in selected}:
            continue
        selected.append(image)

    return main_image, selected[:10]


def generate_brochure_copy(
    client: OpenAI,
    model: str,
    listing: ListingExtraction,
    main_image: SelectedImage,
    gallery_images: Sequence[SelectedImage],
) -> BrochureCopyResponse:
    gallery_summary = "\n".join(
        f"- {image.path.name} | scene={image.scene_type} | room={image.room_type} | draft={image.caption}"
        for image in gallery_images
    )
    response = client.responses.parse(
        model=model,
        instructions=(
            "Write brochure copy for a real-estate flyer. "
            "Return one title, 6-7 summary bullets, and one caption object for each gallery image. "
            "The title must read like a polished headline, not a list of features. "
            "It should evoke the home and lifestyle, be specific, and stay in title case. "
            "Avoid comma-spliced feature lists. "
            "Each summary bullet should be a complete, brochure-ready point of roughly 18-32 words. "
            "Each caption object must include the exact file_name provided for that image. "
            "Each caption should be balanced and concise, ideally 70-110 characters, and should fit a small caption box. "
            "Do not repeat nearly identical captions. "
            "Do not mention file names in the caption text."
        ),
        input=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_text",
                        "text": (
                            f"Short address: {listing.short_address}\n"
                            f"Full address: {listing.address}\n"
                            f"Price: {listing.price}\n"
                            f"Beds/Baths: {listing.bedrooms} / {listing.bathrooms}\n"
                            f"Features: {format_feature_rows_for_prompt(listing.feature_rows)}\n"
                            f"Listing bullets: {' | '.join(listing.summary_bullets)}\n"
                            f"Main image: scene={main_image.scene_type} draft={main_image.caption}\n"
                            f"Other images in required caption order:\n{gallery_summary}"
                        ),
                    }
                ],
            }
        ],
        text_format=BrochureCopyResponse,
        text={"verbosity": "low"},
        reasoning={"effort": "low"},
        max_output_tokens=600,
    )
    parsed = response.output_parsed
    if not parsed or not parsed.title:
        raise RuntimeError("The model returned no brochure copy.")
    parsed.title = " ".join(parsed.title.split())
    parsed.summary_bullets = [clean_bullet_text(item) for item in parsed.summary_bullets if item.strip()]
    parsed.captions = [
        CaptionItem(file_name=item.file_name, caption=clean_caption(item.caption))
        for item in parsed.captions
        if item.file_name.strip() and item.caption.strip()
    ]
    return parsed


def format_feature_rows_for_prompt(feature_rows: Sequence[FeatureRow]) -> str:
    parts: list[str] = []
    for row in feature_rows:
        values = ", ".join(value.strip() for value in row.values if value.strip())
        if values:
            parts.append(f"{row.label}: {values}")
    return " | ".join(parts)


def clean_bullet_text(text: str) -> str:
    return " ".join(text.split()).lstrip("-• ").strip()


def is_size_line(text: str) -> bool:
    value = " ".join(text.split()).strip()
    if not value:
        return False
    if value == "x":
        return True
    return bool(re.search(r"\d+'\s*\d*\"?\s*x\s*\d+'\s*\d*\"?", value))


def build_text_mapping(
    listing: ListingExtraction, brochure_copy: BrochureCopyResponse, gallery_images: Sequence[SelectedImage]
) -> dict[str, str]:
    caption_map = {item.file_name: item.caption for item in brochure_copy.captions}
    mapping = {
        "{{ADDRESS}}": listing.short_address,
        "{{PRICE}}": listing.price,
        "{{BEDS_AND_BATHS}}": format_beds_and_baths(listing.bedrooms, listing.bathrooms),
        "{{TITLE}}": brochure_copy.title,
        "{{TOTAL}}": format_total(listing.total_sqft),
    }
    for index in range(1, 11):
        caption = ""
        if index <= len(gallery_images):
            image = gallery_images[index - 1]
            caption = caption_map.get(image.path.name, image.caption)
        mapping[f"{{{{caption_{index}}}}}"] = caption
    return mapping


def format_beds_and_baths(bedrooms: int | None, bathrooms: int | None) -> str:
    bed_text = f"{bedrooms} Bedroom" if bedrooms == 1 else f"{bedrooms} Bedrooms" if bedrooms else "Bedrooms"
    bath_text = f"{bathrooms} Bath" if bathrooms == 1 else f"{bathrooms} Baths" if bathrooms else "Baths"
    return f"{bed_text} & {bath_text}"


def format_total(total_sqft: str | None) -> str:
    if not total_sqft:
        return "Total"
    normalized = normalize_sqft_text(total_sqft)
    return f"Total ({normalized})"


def normalize_sqft_text(value: str) -> str:
    compact = " ".join(value.split())
    compact = re.sub(r"\(Total\)", "", compact, flags=re.IGNORECASE).strip()
    match = re.search(r"[\d,.]+", compact)
    if not match:
        return compact
    return f"{match.group(0)} sq. ft."


def extract_area_hints(raw_text: str) -> dict[str, str]:
    hints: dict[str, str] = {}
    patterns = {
        "main": [r"Above Grade\s+([\d,.]+)", r"1st\s+([\d,.]+)"],
        "lower": [r"Below Grade\s+([\d,.]+)", r"Bsmt\s+([\d,.]+)"],
    }
    for key, regexes in patterns.items():
        for pattern in regexes:
            match = re.search(pattern, raw_text, flags=re.IGNORECASE)
            if match:
                hints[key] = normalize_sqft_text(match.group(1))
                break
    return hints


def clean_feature_rows(feature_rows: Sequence[FeatureRow], raw_text: str) -> list[FeatureRow]:
    values_by_label = {row.label.strip().lower(): [v.strip() for v in row.values if v.strip()] for row in feature_rows}

    location_values = values_by_label.get("location", [])
    location = ""
    for value in location_values:
        if "-" in value:
            candidate = value.split("-", 1)[1].strip()
            if candidate:
                location = candidate
                break
        if value and value.lower() != "kamloops and district":
            location = value
            break
    if not location:
        for area in ["Brocklehurst", "Aberdeen", "Sahali", "Valleyview", "Juniper Ridge", "Westsyde", "North Kamloops"]:
            if re.search(area, raw_text, flags=re.IGNORECASE):
                location = area
                break

    style = next((value for value in values_by_label.get("style", []) if len(value) < 40), "")

    heating_clean: list[str] = []
    for candidate in ["Forced Air", "Central Air", "Natural Gas"]:
        if any(candidate.lower() in value.lower() for value in values_by_label.get("heating", [])):
            heating_clean.append(candidate)

    parking_source = " | ".join(values_by_label.get("parking", []))
    parking_clean: list[str] = []
    garage_match = re.search(r"Garage Spaces[: ]+(\d+(?:\.\d+)?)", parking_source, flags=re.IGNORECASE)
    total_match = re.search(r"Parking Total[: ]+(\d+(?:\.\d+)?)", parking_source, flags=re.IGNORECASE)
    if garage_match:
        garage_spaces = int(float(garage_match.group(1)))
        if garage_spaces > 0:
            parking_clean.append(f"{garage_spaces} Car Garage")
    if total_match and garage_match and int(float(total_match.group(1))) > int(float(garage_match.group(1))):
        parking_clean.append("Additional Parking")
    if not parking_clean:
        garage_match = re.search(r"Garage Spaces\s+(\d+(?:\.\d+)?)", raw_text, flags=re.IGNORECASE)
        total_match = re.search(r"Parking Total\s+(\d+(?:\.\d+)?)", raw_text, flags=re.IGNORECASE)
        if garage_match:
            garage_spaces = int(float(garage_match.group(1)))
            if garage_spaces > 0:
                parking_clean.append(f"{garage_spaces} Car Garage")
        if total_match and garage_match and int(float(total_match.group(1))) > int(float(garage_match.group(1))):
            parking_clean.append("Additional Parking")

    include_clean: list[str] = []
    include_source = values_by_label.get("includes", [])
    for appliance in ["Dishwasher", "Electric Range", "Refrigerator", "Washer/Dryer"]:
        if any(appliance.lower() in value.lower() for value in include_source):
            include_clean.append(appliance)
    for extra_feature, patterns in {
        "Deck": [r"\bDeck\b", r"\bCovered, Deck\b"],
        "Balcony": [r"\bBalcony\b"],
        "Private Yard": [r"\bPrivate Yard\b", r"\bfenced backyard\b", r"\bfenced yard\b"],
    }.items():
        if any(re.search(pattern, raw_text, flags=re.IGNORECASE) for pattern in patterns):
            include_clean.append(extra_feature)
    include_clean = list(dict.fromkeys(include_clean))

    age_clean: list[str] = []
    age_source = " | ".join(values_by_label.get("age", []))
    built_match = re.search(r"(?:Year Built|Built)\s*(\d{4})", age_source, flags=re.IGNORECASE)
    if not built_match:
        built_match = re.search(r"Year Built\s+(\d{4})", raw_text, flags=re.IGNORECASE)
    if built_match:
        built_year = int(built_match.group(1))
        listing_year_match = re.search(r"Date Listed\s+\w+\s+\d{1,2}/(\d{2})", raw_text)
        listing_year = 2000 + int(listing_year_match.group(1)) if listing_year_match else 2025
        age_clean.append(f"{listing_year - built_year} Years")

    taxes = next((value for value in values_by_label.get("taxes", []) if "$" in value), "")
    lot_size = next((value for value in values_by_label.get("lot size", []) if value), "")

    cleaned: list[FeatureRow] = []
    for label, values in [
        ("Location", [location] if location else []),
        ("Style", [style] if style else []),
        ("Heating", heating_clean),
        ("Parking", parking_clean),
        ("Includes", include_clean),
        ("Age", age_clean),
        ("Taxes", [taxes] if taxes else []),
        ("Lot Size", [lot_size] if lot_size else []),
    ]:
        if values:
            cleaned.append(FeatureRow(label=label, values=values))
    return cleaned


def clean_room_sections(room_sections: Sequence[RoomSection], raw_text: str) -> list[RoomSection]:
    parsed_from_raw = parse_room_sections_from_raw(raw_text)
    if parsed_from_raw:
        return parsed_from_raw

    area_hints = extract_area_hints(raw_text)
    cleaned_sections: list[RoomSection] = []
    for section in room_sections:
        title = section.title.strip()
        title_lower = title.lower()
        if title_lower in {"1st", "first", "main floor", "upper main"}:
            title = "Main"
        elif title_lower in {"lower floor", "basement", "bsmt", "lower"}:
            title = "Basement"
        elif title_lower in {"2nd", "second", "upstairs", "upper"}:
            title = "Upstairs"

        area_text = section.area_text.strip()
        if not area_text or "not specified" in area_text.lower():
            if title.lower() == "main" and "main" in area_hints:
                area_text = area_hints["main"]
            elif title.lower() == "basement" and "lower" in area_hints:
                area_text = area_hints["lower"]

        rooms: list[RoomEntry] = []
        for room in section.rooms:
            name = room.name.strip()
            size = " ".join(room.size.split()).strip()
            if "Full 4 PCE" in name:
                name = name.replace(" - Full 4 PCE", "")
                size = "4 piece"
            if "Half 2 PCE" in name:
                name = name.replace(" - Half 2 PCE", "")
                size = "2 piece"
            size = size.replace("  ", " ")
            if size == "x":
                size = ""
            rooms.append(RoomEntry(name=name, size=size))
        cleaned_sections.append(RoomSection(title=title, area_text=area_text, rooms=rooms))
    return cleaned_sections


def parse_room_sections_from_raw(raw_text: str) -> list[RoomSection]:
    lines = [line.strip() for line in raw_text.splitlines()]
    try:
        start = next(i for i, line in enumerate(lines) if line.startswith("ROOMS (Total:"))
        end = next(i for i, line in enumerate(lines[start + 1 :], start + 1) if line == "BUILDING")
    except StopIteration:
        return []

    block = lines[start:end]
    area_hints = extract_area_hints(raw_text)

    try:
        total_area_idx = block.index("Total Area")
        bsmt_idx = block.index("Bsmt")
    except ValueError:
        return []

    room_names = [line for line in block[total_area_idx + 1 : bsmt_idx] if line]
    if not room_names:
        return []

    main_sizes: list[str] = []
    basement_sizes: list[str] = []
    trailing_sizes: list[str] = []

    for idx, line in enumerate(block):
        if line == "1st":
            for candidate in block[idx + 1 :]:
                if candidate in {"Beds Suite", "Baths Suite", "BUILDING"}:
                    break
                if is_size_line(candidate):
                    main_sizes.append(candidate)
        if line == "Lower":
            for candidate in block[idx + 1 :]:
                if candidate == "1st":
                    break
                if is_size_line(candidate):
                    basement_sizes.append(candidate)

    # Some basement room sizes spill just after the BUILDING header in the PDF text extraction.
    for candidate in lines[end : min(len(lines), end + 30)]:
        if candidate == "Heating":
            break
        if is_size_line(candidate):
            trailing_sizes.append(candidate)

    basement_sizes.extend(trailing_sizes[: max(0, len(room_names) - len(main_sizes) - len(basement_sizes))])
    if not main_sizes:
        return []

    main_names = room_names[: len(main_sizes)]
    basement_names = room_names[len(main_sizes) :]

    # If one generic bedroom belongs in the basement, move it across until counts match.
    while len(basement_names) < len(basement_sizes):
        moved = False
        for idx in range(len(main_names) - 1, -1, -1):
            if main_names[idx] == "Bedroom":
                basement_names.insert(0, main_names.pop(idx))
                moved = True
                break
        if not moved:
            break

    def normalize_room(name: str, size: str) -> RoomEntry:
        clean_name = name.strip()
        clean_size = " ".join(size.split()).strip()
        if "Full 4 PCE" in clean_name:
            clean_name = "Bathroom"
            clean_size = "4 piece"
        elif "Half 2 PCE" in clean_name:
            clean_name = "Ensuite"
            clean_size = "2 piece"
        elif clean_size == "x":
            clean_size = ""
        return RoomEntry(name=clean_name, size=clean_size)

    sections: list[RoomSection] = []
    if main_names:
        sections.append(
            RoomSection(
                title="Main",
                area_text=area_hints.get("main", ""),
                rooms=[normalize_room(name, size) for name, size in zip(main_names, main_sizes)],
            )
        )
    if basement_names:
        sections.append(
            RoomSection(
                title="Basement",
                area_text=area_hints.get("lower", ""),
                rooms=[normalize_room(name, size) for name, size in zip(basement_names, basement_sizes)],
            )
        )
    return sections


def clone_run_style(source_run, dest_run) -> None:
    if source_run is None:
        return
    font = source_run.font
    dest_font = dest_run.font
    dest_font.bold = font.bold
    dest_font.italic = font.italic
    dest_font.name = font.name
    dest_font.size = font.size
    if font.color is not None and getattr(font.color, "rgb", None) is not None:
        dest_font.color.rgb = font.color.rgb


def apply_paragraph_style(source, dest) -> None:
    dest.alignment = source.alignment
    dest.level = source.level
    dest.line_spacing = source.line_spacing
    dest.space_before = source.space_before
    dest.space_after = source.space_after
    source_ppr = getattr(source._p, "pPr", None)
    if source_ppr is not None:
        dest_ppr = getattr(dest._p, "pPr", None)
        if dest_ppr is not None:
            dest._p.remove(dest_ppr)
        dest._p.insert(0, deepcopy(source_ppr))


def apply_run_style_from_snapshot(run, snapshot: dict) -> None:
    font = run.font
    font.bold = snapshot.get("bold")
    font.italic = snapshot.get("italic")
    font.name = snapshot.get("name")
    font.size = snapshot.get("size")
    color = snapshot.get("color")
    if color is not None:
        font.color.rgb = color


def capture_first_run_style(paragraph) -> dict:
    run = paragraph.runs[0] if paragraph.runs else None
    if run is None:
        return {}
    color = None
    if run.font.color is not None and getattr(run.font.color, "rgb", None) is not None:
        color = run.font.color.rgb
    return {
        "bold": run.font.bold,
        "italic": run.font.italic,
        "name": run.font.name,
        "size": run.font.size,
        "color": color,
    }


def replace_text_in_shape(shape, replacements: dict[str, str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        original_text = "".join(run.text for run in paragraph.runs)
        if not original_text:
            continue
        updated_text = original_text
        for placeholder, replacement in replacements.items():
            updated_text = updated_text.replace(placeholder, replacement)
        if updated_text == original_text:
            continue

        source_run = paragraph.runs[0] if paragraph.runs else None
        alignment = paragraph.alignment
        level = paragraph.level
        paragraph.text = updated_text
        paragraph.alignment = alignment
        paragraph.level = level
        if paragraph.runs:
            clone_run_style(source_run, paragraph.runs[0])


def populate_summary_block(shape, bullets: Sequence[str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    paragraph_style = template_paragraph
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()

    values = list(bullets)[:7] or [""]
    for index, bullet in enumerate(values):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        apply_paragraph_style(paragraph_style, paragraph)
        paragraph.text = bullet
        if paragraph.runs:
            apply_run_style_from_snapshot(paragraph.runs[0], run_style)
    distribute_summary_spacing(shape, text_frame, run_style)


def distribute_summary_spacing(shape, text_frame, run_style: dict) -> None:
    paragraphs = [p for p in text_frame.paragraphs if "".join(run.text for run in p.runs).strip()]
    if len(paragraphs) < 2:
        return

    font_size = run_style.get("size")
    font_points = font_size.pt if font_size is not None else 14
    estimated_used = len(paragraphs) * font_points * 1.8
    remaining = max(0.0, shape.height.pt - estimated_used)
    per_gap = min(6.0, remaining / max(1, len(paragraphs) - 1))

    first_seen = False
    for paragraph in paragraphs:
        if not first_seen:
            paragraph.space_before = Pt(0)
            first_seen = True
        else:
            paragraph.space_before = Pt(per_gap)
        paragraph.space_after = Pt(0)


def populate_feature_block(shape, feature_rows: Sequence[FeatureRow]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    paragraph_style = template_paragraph
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()
    label_count = 0
    child_count = 0

    first = True
    for row in feature_rows:
        values = [value.strip() for value in row.values if value.strip()]
        if not values:
            continue

        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        label_count += 1
        apply_paragraph_style(paragraph_style, paragraph)
        paragraph.text = f"{row.label}:\t{values[0]}"
        for run in paragraph.runs:
            apply_run_style_from_snapshot(run, run_style)

        for extra_value in values[1:]:
            extra = text_frame.add_paragraph()
            child_count += 1
            apply_paragraph_style(paragraph_style, extra)
            extra.text = f"\t{extra_value}"
            for run in extra.runs:
                apply_run_style_from_snapshot(run, run_style)

    distribute_feature_spacing(shape, text_frame, label_count, child_count, run_style)


def distribute_feature_spacing(shape, text_frame, label_count: int, child_count: int, run_style: dict) -> None:
    paragraphs = [p for p in text_frame.paragraphs if "".join(run.text for run in p.runs).strip()]
    if not paragraphs:
        return

    font_size = run_style.get("size")
    font_points = font_size.pt if font_size is not None else 14
    base_line_points = font_points * 1.2
    estimated_used = len(paragraphs) * base_line_points
    remaining = max(0.0, shape.height.pt - estimated_used)
    total_units = max(1, label_count * 2 + child_count)
    points_per_unit = min(3.5, remaining / total_units)

    first_seen = False
    for paragraph in paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not first_seen:
            paragraph.space_before = Pt(0)
            first_seen = True
            continue
        units = 1
        if ":\t" in text:
            units = 2
        paragraph.space_before = Pt(points_per_unit * units)
        paragraph.space_after = Pt(0)


def populate_room_block(shape, room_sections: Sequence[RoomSection]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    from pptx.enum.text import PP_ALIGN

    text_frame = shape.text_frame
    template_paragraph = text_frame.paragraphs[0]
    paragraph_style = template_paragraph
    run_style = capture_first_run_style(template_paragraph)
    text_frame.clear()
    first = True

    for section in room_sections:
        heading = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        apply_paragraph_style(paragraph_style, heading)
        heading.alignment = PP_ALIGN.CENTER
        heading.text = f"{section.title} ({normalize_sqft_text(section.area_text)})"
        for run in heading.runs:
            apply_run_style_from_snapshot(run, run_style)
            run.font.italic = True

        spacer = text_frame.add_paragraph()
        spacer.text = ""

        for room in section.rooms:
            row = text_frame.add_paragraph()
            apply_paragraph_style(paragraph_style, row)
            row.alignment = None
            row.text = f"{room.name}\t{room.size}"
            for run in row.runs:
                apply_run_style_from_snapshot(run, run_style)

def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def fill_picture_placeholder(slide, placeholder_shape, image_path: Path) -> None:
    from pptx.util import Emu

    left, top, width, height = (
        placeholder_shape.left,
        placeholder_shape.top,
        placeholder_shape.width,
        placeholder_shape.height,
    )

    with Image.open(image_path) as img:
        image_width, image_height = img.size

    image_ratio = image_width / image_height
    box_ratio = width / height

    picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    if image_ratio > box_ratio:
        scaled_width = height * image_ratio
        crop = (scaled_width - width) / scaled_width / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        scaled_height = width / image_ratio
        crop = (scaled_height - height) / scaled_height / 2
        picture.crop_top = crop
        picture.crop_bottom = crop

    remove_shape(placeholder_shape)


def render_presentation(
    template_path: Path,
    output_path: Path,
    text_mapping: dict[str, str],
    feature_rows: Sequence[FeatureRow],
    room_sections: Sequence[RoomSection],
    main_image: SelectedImage,
    gallery_images: Sequence[SelectedImage],
) -> None:
    pptx = require_dependency("pptx", "python-pptx")
    presentation = pptx.Presentation(str(template_path))

    image_map = {"{{IMG_MAIN}}": main_image.path}
    for index, image in enumerate(gallery_images, start=1):
        image_map[f"{{{{IMG_{index}}}}}"] = image.path

    for slide in presentation.slides:
        for shape in list(slide.shapes):
            shape_name = getattr(shape, "name", "").strip()
            if shape_name in image_map:
                fill_picture_placeholder(slide, shape, image_map[shape_name])
                continue

            if getattr(shape, "has_text_frame", False):
                combined_text = "\n".join("".join(run.text for run in paragraph.runs) for paragraph in shape.text_frame.paragraphs)
                if "{{SUMMARY}}" in combined_text:
                    populate_summary_block(shape, text_mapping.get("{{SUMMARY}}", "").split("\n"))
                    continue
                if "{{FEATURE_BLOCK}}" in combined_text:
                    populate_feature_block(shape, feature_rows)
                    continue
                if "{{ROOM_BLOCK}}" in combined_text:
                    populate_room_block(shape, room_sections)
                    continue
                replace_text_in_shape(shape, text_mapping)

    presentation.save(str(output_path))


def run_pipeline(config: AppConfig, logger: Callable[[str], None] = log_print) -> Path:
    logger("Reading listing PDF...")
    raw_text = extract_pdf_text(config.pdf_path)

    client = build_client()

    logger("Extracting brochure fields from listing text...")
    listing = analyze_listing_text(client, config.model, raw_text)
    listing.feature_rows = clean_feature_rows(listing.feature_rows, raw_text)
    listing.room_sections = clean_room_sections(listing.room_sections, raw_text)

    logger("Scanning image folder...")
    image_paths = iter_image_paths(config.image_dir)

    logger(f"Found {len(image_paths)} image(s).")
    ranked_images = analyze_images(client, config.model, image_paths, logger)
    main_image, gallery_images = choose_images(ranked_images)

    logger("Generating brochure copy...")
    brochure_copy = generate_brochure_copy(client, config.model, listing, main_image, gallery_images)
    caption_map = {item.file_name: item.caption for item in brochure_copy.captions}
    for image in gallery_images:
        if image.path.name in caption_map:
            image.caption = caption_map[image.path.name]

    text_mapping = build_text_mapping(listing, brochure_copy, gallery_images)
    text_mapping["{{SUMMARY}}"] = "\n".join(brochure_copy.summary_bullets or listing.summary_bullets)

    logger("Rendering PowerPoint brochure...")
    render_presentation(
        template_path=config.template_path,
        output_path=config.output_path,
        text_mapping=text_mapping,
        feature_rows=listing.feature_rows,
        room_sections=listing.room_sections,
        main_image=main_image,
        gallery_images=gallery_images,
    )

    logger(f"Saved brochure to {config.output_path}")
    return config.output_path


def parse_args(argv: Sequence[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a real-estate brochure PPTX from a listing PDF and image folder.")
    parser.add_argument("--pdf", type=Path, help="Path to the listing PDF.")
    parser.add_argument("--images", type=Path, help="Path to the image directory.")
    parser.add_argument("--template", type=Path, default=Path("Brochure_Template.pptx"), help="Path to the PPTX template.")
    parser.add_argument("--output", type=Path, default=Path("brochure_output.pptx"), help="Output PPTX path.")
    parser.add_argument("--model", default=DEFAULT_MODEL, help="OpenAI model name.")
    parser.add_argument("--cli", action="store_true", help="Force CLI mode even if tkinter is available.")
    return parser.parse_args(argv)


def validate_config(config: AppConfig) -> None:
    if not config.pdf_path.exists():
        raise RuntimeError(f"PDF not found: {config.pdf_path}")
    if not config.image_dir.exists() or not config.image_dir.is_dir():
        raise RuntimeError(f"Image directory not found: {config.image_dir}")
    if not config.template_path.exists():
        raise RuntimeError(f"Template not found: {config.template_path}")
    config.output_path.parent.mkdir(parents=True, exist_ok=True)


def maybe_launch_gui(default_template: Path, default_model: str) -> bool:
    try:
        import tkinter as tk
        from tkinter import filedialog, messagebox
        from tkinter.scrolledtext import ScrolledText
    except Exception:
        return False

    class BrochureApp:
        def __init__(self, root: tk.Tk) -> None:
            self.root = root
            self.root.title("Brochure Agent")
            self.root.geometry("760x540")

            self.pdf_var = tk.StringVar()
            self.images_var = tk.StringVar()
            self.template_var = tk.StringVar(value=str(default_template))
            self.output_var = tk.StringVar(value=str(Path.cwd() / "brochure_output.pptx"))
            self.model_var = tk.StringVar(value=default_model)

            self._build()

        def _build(self) -> None:
            frame = tk.Frame(self.root, padx=12, pady=12)
            frame.pack(fill="both", expand=True)

            self._row(frame, 0, "Listing PDF", self.pdf_var, lambda: self.pdf_var.set(filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])))
            self._row(frame, 1, "Image Folder", self.images_var, lambda: self.images_var.set(filedialog.askdirectory()))
            self._row(frame, 2, "Template", self.template_var, lambda: self.template_var.set(filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])))
            self._row(frame, 3, "Output", self.output_var, lambda: self.output_var.set(filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])))
            self._row(frame, 4, "Model", self.model_var, None)

            button_frame = tk.Frame(frame)
            button_frame.grid(row=5, column=0, columnspan=3, sticky="ew", pady=(10, 10))
            tk.Button(button_frame, text="Generate Brochure", command=self.start_generation).pack(side="left")

            self.log_widget = ScrolledText(frame, height=20, wrap="word")
            self.log_widget.grid(row=6, column=0, columnspan=3, sticky="nsew")

            frame.columnconfigure(1, weight=1)
            frame.rowconfigure(6, weight=1)

        def _row(self, parent, row: int, label: str, variable, browse_command) -> None:
            tk.Label(parent, text=label, anchor="w").grid(row=row, column=0, sticky="w", pady=4)
            tk.Entry(parent, textvariable=variable).grid(row=row, column=1, sticky="ew", padx=(8, 8))
            if browse_command:
                tk.Button(parent, text="Browse", command=browse_command).grid(row=row, column=2, sticky="ew")

        def append_log(self, message: str) -> None:
            self.log_widget.insert("end", message + "\n")
            self.log_widget.see("end")
            self.root.update_idletasks()

        def start_generation(self) -> None:
            config = AppConfig(
                pdf_path=Path(self.pdf_var.get()),
                image_dir=Path(self.images_var.get()),
                template_path=Path(self.template_var.get()),
                output_path=Path(self.output_var.get()),
                model=self.model_var.get().strip() or DEFAULT_MODEL,
            )
            try:
                validate_config(config)
            except Exception as exc:
                messagebox.showerror("Invalid input", str(exc))
                return

            thread = threading.Thread(target=self._run_generation, args=(config,), daemon=True)
            thread.start()

        def _run_generation(self, config: AppConfig) -> None:
            try:
                run_pipeline(config, logger=self.append_log)
                messagebox.showinfo("Success", f"Saved brochure to:\n{config.output_path}")
            except Exception as exc:
                self.append_log(traceback.format_exc())
                messagebox.showerror("Generation failed", str(exc))

    root = tk.Tk()
    BrochureApp(root)
    root.mainloop()
    return True


def run_cli(args: argparse.Namespace) -> int:
    if not args.pdf or not args.images:
        print("CLI mode requires --pdf and --images.", file=sys.stderr)
        return 2

    config = AppConfig(
        pdf_path=args.pdf,
        image_dir=args.images,
        template_path=args.template,
        output_path=args.output,
        model=args.model,
    )
    try:
        validate_config(config)
        run_pipeline(config)
        return 0
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1


def main(argv: Sequence[str] | None = None) -> int:
    args = parse_args(argv or sys.argv[1:])
    if args.cli:
        return run_cli(args)

    if maybe_launch_gui(args.template, args.model):
        return 0
    return run_cli(args)


if __name__ == "__main__":
    raise SystemExit(main())
